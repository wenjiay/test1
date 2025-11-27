from pptx import Presentation
import fitz  # PyMuPDF

def extract_text_from_file(filepath):
    if filepath.endswith('.pdf'):
        doc = fitz.open(filepath)
        text = ""
        for page in doc:
            text += page.get_text()
        return text

    elif filepath.endswith('.pptx'):
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    else:
        return None
